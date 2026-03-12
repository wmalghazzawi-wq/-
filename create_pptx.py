from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_FILE = Path(__file__).resolve().parent / "nanotechnology_website_summary.pptx"

# Website theme colors from styles.css
BG_DARK = RGBColor(5, 11, 24)          # #050b18
BG_LAYER = RGBColor(10, 22, 40)        # close to #0a1628 in hero gradient
ACCENT_CYAN = RGBColor(0, 247, 255)    # #00f7ff
ACCENT_BLUE = RGBColor(59, 130, 246)   # #3b82f6
ACCENT_ORANGE = RGBColor(255, 107, 53) # #ff6b35
TEXT_PRIMARY = RGBColor(232, 244, 253) # #e8f4fd
TEXT_SECONDARY = RGBColor(139, 163, 204)  # #8ba3cc

FONT_TITLE = "Arial"
FONT_BODY = "Arial"


def apply_background(slide, section_tag: str | None = None) -> None:
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = BG_DARK

    top_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32)
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = BG_LAYER
    top_band.line.color.rgb = BG_LAYER

    accent_line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.30), Inches(13.333), Inches(0.05)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT_BLUE
    accent_line.line.color.rgb = ACCENT_BLUE

    glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.8), Inches(5.8), Inches(3.6), Inches(3.6)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = ACCENT_CYAN
    glow.line.color.rgb = ACCENT_CYAN

    if section_tag:
        tag = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(10.55),
            Inches(0.48),
            Inches(2.45),
            Inches(0.48),
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = ACCENT_ORANGE
        tag.line.color.rgb = ACCENT_ORANGE
        tag_text = tag.text_frame
        tag_text.text = section_tag
        p = tag_text.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_BODY
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)


def add_footer(slide, text: str = "موقع عالم النانو") -> None:
    box = slide.shapes.add_textbox(Inches(0.65), Inches(7.05), Inches(12.0), Inches(0.3))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_SECONDARY


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, "ملخص الموقع")

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(11.8), Inches(1.25))
    tf = title_box.text_frame
    tf.text = "عالم النانو"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT_TITLE
    p.font.bold = True
    p.font.size = Pt(52)
    p.font.color.rgb = ACCENT_CYAN

    subtitle_box = slide.shapes.add_textbox(Inches(0.82), Inches(2.75), Inches(11.6), Inches(1.5))
    stf = subtitle_box.text_frame
    stf.text = "ملخص شامل لمحتوى الموقع، أقسامه التعليمية، ونظام التنبؤ الذكي"
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.RIGHT
    sp.font.name = FONT_BODY
    sp.font.size = Pt(24)
    sp.font.color.rgb = TEXT_PRIMARY

    chips_box = slide.shapes.add_textbox(Inches(0.84), Inches(4.15), Inches(11.6), Inches(1.5))
    ctf = chips_box.text_frame
    ctf.clear()
    for line in [
        "• المواد النانوية", "• التطبيقات", "• التنبؤ بالصور", "• تحليل XRD"
    ]:
        p = ctf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = FONT_BODY
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_SECONDARY

    add_footer(slide, "الهوية البصرية: نفس خلفية وألوان الموقع")


def add_bullet_slide(
    prs: Presentation,
    heading: str,
    bullets: list[str],
    section_tag: str,
    accent_color: RGBColor = ACCENT_CYAN,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, section_tag)

    heading_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.6), Inches(0.9))
    htf = heading_box.text_frame
    htf.text = heading
    hp = htf.paragraphs[0]
    hp.alignment = PP_ALIGN.RIGHT
    hp.font.name = FONT_TITLE
    hp.font.bold = True
    hp.font.size = Pt(36)
    hp.font.color.rgb = accent_color

    body_box = slide.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(11.0), Inches(4.6))
    btf = body_box.text_frame
    btf.word_wrap = True
    btf.clear()

    for idx, line in enumerate(bullets):
        p = btf.add_paragraph()
        p.text = f"• {line}"
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = FONT_BODY
        p.font.size = Pt(23 if idx == 0 else 21)
        p.font.color.rgb = TEXT_PRIMARY if idx < 2 else TEXT_SECONDARY
        p.level = 0

    marker = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.DONUT, Inches(0.3), Inches(2.15), Inches(0.4), Inches(0.4)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent_color
    marker.line.color.rgb = accent_color

    add_footer(slide)


def create_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_bullet_slide(
        prs,
        "فكرة الموقع ومحتواه",
        [
            "منصة عربية تعليمية بعنوان عالم النانو لعرض أساسيات تقنية النانو.",
            "تشرح المفاهيم من 1 إلى 100 نانومتر بلغة مبسطة ومدعومة بصريًا.",
            "تضم أقسامًا رئيسية: المواد النانوية، التطبيقات، وحول المشروع.",
            "تعرض حقائق علمية سريعة وروابط تنقل واضحة بين الصفحات.",
        ],
        "الرؤية",
    )

    add_bullet_slide(
        prs,
        "أبرز المواد النانوية في الموقع",
        [
            "أنابيب الكربون النانوية، أكاسيد المعادن، الجرافين، النقاط الكمومية.",
            "يشمل أيضًا الدندريمرات والأسلاك النانوية مع استخدامات كل نوع.",
            "أكسيد الزنك ZnO وأكسيد النيكل NiO ضمن سياق التحليل والتصنيف.",
            "تنظيم المحتوى بشكل بطاقات يسهل المقارنة والفهم.",
        ],
        "المواد",
        accent_color=ACCENT_BLUE,
    )

    add_bullet_slide(
        prs,
        "تطبيقات تقنية النانو",
        [
            "الطب والرعاية الصحية: توصيل دوائي، استشعار حيوي، وتحسين التشخيص.",
            "الطاقة والبيئة: معالجة المياه وتحسين أداء الخلايا الشمسية.",
            "الإلكترونيات والحوسبة: مواد متقدمة لترانزستورات الجيل القادم.",
            "المواد والطلاءات: أسطح مضادة للبكتيريا وطلاءات ذاتية التنظيف.",
        ],
        "التطبيقات",
        accent_color=ACCENT_ORANGE,
    )

    add_bullet_slide(
        prs,
        "نظام التنبؤ الذكي بالصور",
        [
            "رفع صورة العينة ثم الضغط على زر ابدأ التنبؤ لإظهار الفئة المتوقعة.",
            "النظام يعرض نوع العينة مع نسبة الثقة وأعلى الاحتمالات البديلة.",
            "يدعم تحليل عينات ZnO وNiO ضمن واجهة تفاعلية مباشرة.",
            "يعتمد على نموذج تعلم عميق مدرب ومربوط بخادم Flask محلي.",
        ],
        "الذكاء الاصطناعي",
    )

    add_bullet_slide(
        prs,
        "أهمية XRD في تحليل العينات النانوية",
        [
            "XRD يحدد البنية البلورية ويميز الأطوار بدقة علمية عالية.",
            "يُستخدم لتقدير حجم البلورات النانوية من قمم وأنماط الحيود.",
            "يساعد على ربط خصائص المادة بالأداء الوظيفي في التطبيقات.",
            "دمج XRD مع الذكاء الاصطناعي يسرّع التصنيف ويقلل الخطأ البشري.",
        ],
        "XRD",
        accent_color=ACCENT_BLUE,
    )

    add_bullet_slide(
        prs,
        "خلاصة",
        [
            "الموقع يجمع بين الشرح العلمي والواجهة التفاعلية في منصة واحدة.",
            "الهوية البصرية في العرض مطابقة لألوان وخلفية الموقع.",
            "يمكن توسيع العرض لاحقًا بإضافة نتائج النموذج والرسوم البيانية.",
        ],
        "الخلاصة",
    )

    prs.save(OUTPUT_FILE)
    return OUTPUT_FILE


if __name__ == "__main__":
    out = create_presentation()
    print(f"Presentation saved: {out}")
